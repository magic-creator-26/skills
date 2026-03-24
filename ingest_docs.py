#!/usr/bin/env python3
"""
enterprise-knowledge ingestion script
--------------------------------------
Reads a folder of documents (PDF, DOCX, PPTX) and uses Claude to extract
domain terms, populating the three reference files in the skill:
  - references/entities.md
  - references/alltables.md
  - references/internal-terms.md

Usage:
  python scripts/ingest_docs.py --docs-folder /path/to/docs --skill-dir /path/to/enterprise-knowledge
  python scripts/ingest_docs.py --file /path/to/single.pdf --skill-dir .
  python scripts/ingest_docs.py --docs-folder /path/to/docs --dry-run   # preview only, no writes

Dependencies:
  pip install anthropic pymupdf python-docx python-pptx tqdm
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text_pdf(path: Path) -> str:
    try:
        import fitz  # pymupdf
        doc = fitz.open(str(path))
        return "\n".join(page.get_text() for page in doc)
    except Exception as e:
        print(f"  [WARN] PDF extraction failed for {path.name}: {e}")
        return ""


def extract_text_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        print(f"  [WARN] DOCX extraction failed for {path.name}: {e}")
        return ""


def extract_text_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        print(f"  [WARN] PPTX extraction failed for {path.name}: {e}")
        return ""


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_pdf(path)
    elif suffix in (".docx", ".doc"):
        return extract_text_docx(path)
    elif suffix in (".pptx", ".ppt"):
        return extract_text_pptx(path)
    else:
        return ""


# ---------------------------------------------------------------------------
# Claude extraction
# ---------------------------------------------------------------------------

EXTRACTION_SYSTEM = """
You are a knowledge extraction assistant for an enterprise data platform.
Your job is to read documentation and extract domain terms into a structured glossary.

The organization has three knowledge domains:
1. **entities** — the Entities system: entity types (Person, Organization, Account…),
   their attributes, canonical identifier fields (Phone, NationalID, AccountID, Name…),
   lifecycle states, and ontology concepts.
2. **alltables** — the AllTables system: a canvas tool where users select org-wide tables
   by name and draw column-level connections. Terms here include table types, table names,
   column names, canvas relationships, source system mappings (CRM, Billing, Events…).
3. **internal** — org-specific vocabulary, Hebrew terms (include both script and transliteration),
   abbreviations, platform-specific concepts, and any jargon that isn't self-explanatory.

Output ONLY a JSON array. No prose, no markdown fences. Each item:
{
  "term": "<canonical term name, English or Hebrew+transliteration>",
  "domain": "entities" | "alltables" | "internal",
  "type": "entity_type" | "table" | "column" | "identifier" | "metric" | "concept" | "abbreviation" | "other",
  "hebrew_alias": "<Hebrew script if applicable, else null>",
  "definition": "<clear 1-3 sentence definition>",
  "attributes": ["field1 — description", "field2 — description"],  // only if entity or table
  "notes": "<edge cases, caveats, relationships to other terms, or null>"
}

Rules:
- Extract EVERY named concept, table, column, field, entity type, or internal term you find.
- If a term is ambiguous across domains, emit one entry per domain with a note.
- For Hebrew terms: use the Hebrew script as `hebrew_alias` and a transliteration as `term`.
- Skip generic programming/SQL terms (SELECT, JOIN, NULL, etc.) — only org-specific terms.
- Minimum 5 terms per document, maximum 80. Focus on quality over quantity.
""".strip()


def extract_terms_from_text(text: str, source_filename: str, client) -> list[dict]:
    """Call Claude to extract terms from a chunk of text."""
    # Truncate to ~12k chars to stay within comfortable context
    chunk = text[:12000]
    if len(text) > 12000:
        chunk += f"\n\n[... document truncated, {len(text) - 12000} chars omitted ...]"

    prompt = f"Source document: {source_filename}\n\n{chunk}"

    response = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=4096,
        system=EXTRACTION_SYSTEM,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = response.content[0].text.strip()
    # Strip accidental markdown fences
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    try:
        terms = json.loads(raw)
        if not isinstance(terms, list):
            raise ValueError("Expected JSON array")
        return terms
    except Exception as e:
        print(f"  [WARN] JSON parse failed: {e}")
        print(f"  Raw response snippet: {raw[:200]}")
        return []


# ---------------------------------------------------------------------------
# Reference file management
# ---------------------------------------------------------------------------

INSERTION_MARKER = "<!-- INGESTED ENTRIES BELOW — do not remove this line, used as insertion marker -->"


def load_existing_terms(ref_path: Path) -> set[str]:
    """Return the set of already-indexed term names (lowercased for dedup)."""
    if not ref_path.exists():
        return set()
    content = ref_path.read_text(encoding="utf-8")
    # Match ## headings
    return {m.group(1).strip().lower() for m in re.finditer(r"^## (.+)$", content, re.MULTILINE)}


def format_entry(term: dict) -> str:
    lines = [f"## {term['term']}"]
    lines.append(f"**System**: {term['domain'].capitalize()}")
    lines.append(f"**Type**: {term.get('type', 'concept')}")
    if term.get("hebrew_alias"):
        lines.append(f"**Hebrew/Alias**: {term['hebrew_alias']}")
    lines.append("")
    lines.append(term.get("definition", "").strip())
    attrs = term.get("attributes") or []
    if attrs:
        lines.append("")
        lines.append("**Fields / Attributes**:")
        for a in attrs:
            lines.append(f"- {a}")
    notes = term.get("notes")
    if notes and notes.strip() and notes.strip().lower() != "null":
        lines.append("")
        lines.append(f"**Notes**: {notes.strip()}")
    lines.append("")
    lines.append("---")
    lines.append("")
    return "\n".join(lines)


def append_entries_to_ref(ref_path: Path, new_terms: list[dict], dry_run: bool) -> int:
    """Append non-duplicate entries to a reference file. Returns count added."""
    existing = load_existing_terms(ref_path)
    to_add = [t for t in new_terms if t["term"].strip().lower() not in existing]

    if not to_add:
        return 0

    if dry_run:
        for t in to_add:
            print(f"    [DRY RUN] Would add: {t['term']}")
        return len(to_add)

    content = ref_path.read_text(encoding="utf-8")
    if INSERTION_MARKER not in content:
        content += f"\n\n{INSERTION_MARKER}\n"

    new_block = "\n".join(format_entry(t) for t in to_add)
    content = content.replace(INSERTION_MARKER, f"{INSERTION_MARKER}\n{new_block}")
    ref_path.write_text(content, encoding="utf-8")
    return len(to_add)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Ingest docs into enterprise-knowledge skill")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--docs-folder", type=Path, help="Folder of documents to ingest")
    group.add_argument("--file", type=Path, help="Single document to ingest")
    parser.add_argument("--skill-dir", type=Path, default=Path("."),
                        help="Path to enterprise-knowledge skill directory (default: cwd)")
    parser.add_argument("--dry-run", action="store_true",
                        help="Extract and print terms without writing to files")
    args = parser.parse_args()

    # Validate skill dir
    refs_dir = args.skill_dir / "references"
    if not refs_dir.exists():
        print(f"[ERROR] References directory not found: {refs_dir}")
        print("Make sure --skill-dir points to the enterprise-knowledge folder.")
        sys.exit(1)

    ref_files = {
        "entities": refs_dir / "entities.md",
        "alltables": refs_dir / "alltables.md",
        "internal": refs_dir / "internal-terms.md",
    }

    # Init Anthropic client
    try:
        import anthropic
        client = anthropic.Anthropic()
    except ImportError:
        print("[ERROR] anthropic package not found. Run: pip install anthropic")
        sys.exit(1)

    # Collect files to process
    if args.file:
        files = [args.file] if args.file.exists() else []
    else:
        supported = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}
        files = [f for f in sorted(args.docs_folder.rglob("*")) if f.suffix.lower() in supported]

    if not files:
        print("[WARN] No supported documents found.")
        sys.exit(0)

    print(f"\nEnterprise Knowledge Ingestion")
    print(f"Documents found : {len(files)}")
    print(f"Skill directory : {args.skill_dir.resolve()}")
    print(f"Dry run         : {args.dry_run}")
    print()

    total_added = {"entities": 0, "alltables": 0, "internal": 0}

    for i, doc_path in enumerate(files, 1):
        print(f"[{i}/{len(files)}] {doc_path.name}")

        text = extract_text(doc_path)
        if not text.strip():
            print("  [SKIP] No text extracted.")
            continue

        print(f"  Extracted {len(text):,} chars — calling Claude...")
        terms = extract_terms_from_text(text, doc_path.name, client)
        print(f"  Got {len(terms)} terms")

        # Bucket by domain
        by_domain: dict[str, list] = {"entities": [], "alltables": [], "internal": []}
        for t in terms:
            domain = t.get("domain", "internal")
            if domain not in by_domain:
                domain = "internal"
            by_domain[domain].append(t)

        for domain, domain_terms in by_domain.items():
            if domain_terms:
                added = append_entries_to_ref(ref_files[domain], domain_terms, args.dry_run)
                total_added[domain] += added
                print(f"    {domain:12s}: +{added} new entries ({len(domain_terms)} extracted)")

    print()
    print("=== Summary ===")
    for domain, count in total_added.items():
        print(f"  {domain:12s}: {count} entries added")
    if args.dry_run:
        print("  (dry run — no files written)")
    print()
    print("Done. Reference files updated:")
    for domain, path in ref_files.items():
        if path.exists():
            n = len(load_existing_terms(path))
            print(f"  {path.relative_to(args.skill_dir)}  ({n} total entries)")


if __name__ == "__main__":
    main()
